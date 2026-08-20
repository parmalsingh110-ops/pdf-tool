"""
PDF Tool — FastAPI Backend v2.1 (Security Hardened)
====================================================
High-quality PDF/Office conversions using Python open-source tools.

Security hardening applied:
  - Upload size limits (MAX_UPLOAD_MB env var)
  - Safe error messages (no raw exception leakage)
  - Rate limiting on all expensive endpoints
  - Security response headers
  - CORS restricted to allowed origins
  - DPI clamping, page-count limits
  - LibreOffice concurrency semaphore
  - Safe filenames in all FileResponse
  - /docs disabled in production (ENABLE_DOCS=false)
  - All print() replaced with structured logger
  - iLovePDF API key strictly server-side (env var only)

Run:
    uvicorn main:app --reload --port 8000
"""

import os
import re
import shutil
import tempfile
import subprocess
import uuid
import json
import time
import logging
import asyncio
from pathlib import Path

from fastapi import FastAPI, UploadFile, File, HTTPException, Form, BackgroundTasks, Request, Response
from fastapi.middleware.cors import CORSMiddleware
from fastapi.middleware.gzip import GZipMiddleware
from fastapi.responses import FileResponse, JSONResponse
from starlette.background import BackgroundTask
from starlette.middleware.base import BaseHTTPMiddleware
from dotenv import load_dotenv
from slowapi import Limiter, _rate_limit_exceeded_handler
from slowapi.util import get_remote_address
from slowapi.errors import RateLimitExceeded

load_dotenv()

# ─── Logging Setup ─────────────────────────────────────────────
logging.basicConfig(
    level=logging.INFO,
    format="%(asctime)s | %(levelname)s | %(message)s",
    datefmt="%Y-%m-%dT%H:%M:%S",
)
logger = logging.getLogger("pdf_tool")

# ─── Security Configuration (from environment) ─────────────────
# All limits are configurable via Render Environment Variables.
# Defaults are reasonable for a free-tier public tool.

MAX_UPLOAD_MB: int = int(os.environ.get("MAX_UPLOAD_MB", "50"))
MAX_UPLOAD_BYTES: int = MAX_UPLOAD_MB * 1024 * 1024

MAX_PDF_PAGES: int = int(os.environ.get("MAX_PDF_PAGES", "300"))
MAX_DPI: int = int(os.environ.get("MAX_DPI", "300"))
MIN_DPI: int = 72
MAX_PROCESSING_SECONDS: int = int(os.environ.get("MAX_PROCESSING_SECONDS", "300"))

# Disable /docs and /redoc in production (set ENABLE_DOCS=true to re-enable)
ENABLE_DOCS: bool = os.environ.get("ENABLE_DOCS", "false").lower() == "true"

# CORS: comma-separated list of allowed origins
_raw_origins = os.environ.get(
    "ALLOWED_ORIGINS",
    "https://www.pdfmediasuite.in,https://pdfmediasuite.in,http://localhost:3000,http://localhost:4173,http://127.0.0.1:3000"
)
ALLOWED_ORIGINS: list[str] = [o.strip() for o in _raw_origins.split(",") if o.strip()]

# LibreOffice concurrency: limit concurrent soffice processes to avoid RAM exhaustion
LIBREOFFICE_SEMAPHORE = asyncio.Semaphore(int(os.environ.get("MAX_CONCURRENT_LIBREOFFICE", "2")))

# ─── Rate Limiter Setup ─────────────────────────────────────────
# Uses client IP to track requests.
# X-Forwarded-For header is used behind Render's proxy.
limiter = Limiter(key_func=get_remote_address)

# ─── FastAPI App ─────────────────────────────────────────────────
docs_url = "/docs" if ENABLE_DOCS else None
redoc_url = "/redoc" if ENABLE_DOCS else None
openapi_url = "/openapi.json" if ENABLE_DOCS else None

app = FastAPI(
    title="PDF Tool Backend",
    version="2.1.0",
    docs_url=docs_url,
    redoc_url=redoc_url,
    openapi_url=openapi_url,
)

# Attach rate limiter to app state
app.state.limiter = limiter
app.add_exception_handler(RateLimitExceeded, _rate_limit_exceeded_handler)


# ─── Middleware Stack ───────────────────────────────────────────

# 1. GZip compression
app.add_middleware(GZipMiddleware, minimum_size=500)

# 2. CORS — restricted to known production origins
app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_methods=["GET", "POST", "HEAD", "OPTIONS"],
    allow_headers=["Content-Type", "Accept", "Authorization"],
    expose_headers=["X-Match-Count"],
    allow_credentials=False,
)


# 3. Security Headers Middleware
class SecurityHeadersMiddleware(BaseHTTPMiddleware):
    """Adds security response headers to every response."""

    # CSP for the FastAPI backend responses.
    # The backend only serves JSON / binary file downloads (no HTML pages),
    # so we lock it down very tightly.  The permissive CSP for the React SPA
    # lives in index.html (served by Vite/Render static hosting).
    _CSP = (
        "default-src 'none'; "
        "script-src  'none'; "
        "style-src   'none'; "
        "img-src     'none'; "
        "connect-src 'self'; "
        "frame-ancestors 'none'; "
        "form-action 'self';"
    )

    async def dispatch(self, request: Request, call_next):
        response = await call_next(request)
        response.headers["Content-Security-Policy"] = self._CSP
        response.headers["X-Content-Type-Options"] = "nosniff"
        response.headers["X-Frame-Options"] = "DENY"
        response.headers["Referrer-Policy"] = "strict-origin-when-cross-origin"
        response.headers["Permissions-Policy"] = "geolocation=(), camera=(), microphone=()"
        # HSTS: only send on HTTPS (Render always serves HTTPS in production)
        if request.url.scheme == "https":
            response.headers["Strict-Transport-Security"] = "max-age=31536000; includeSubDomains"
        return response


app.add_middleware(SecurityHeadersMiddleware)


# 4. Bandwidth Logging Middleware — logs IP, route, method, status, response size.
class BandwidthLoggerMiddleware(BaseHTTPMiddleware):
    async def dispatch(self, request: Request, call_next):
        start = time.perf_counter()
        response = await call_next(request)

        # We must NOT buffer the response in memory (bad for large PDFs).
        # Instead, we wrap the streaming iterator to count bytes as they go out.
        original_iterator = response.body_iterator

        async def logging_iterator():
            total_bytes = 0
            try:
                async for chunk in original_iterator:
                    total_bytes += len(chunk)
                    yield chunk
            finally:
                elapsed_ms = round((time.perf_counter() - start) * 1000, 1)
                size_kb = round(total_bytes / 1024, 2)
                # Safely extract client IP — prefer X-Forwarded-For (Render proxy)
                forwarded = request.headers.get("x-forwarded-for", "")
                # Only use the first IP from X-Forwarded-For to prevent IP spoofing
                client_ip = forwarded.split(",")[0].strip() if forwarded else (
                    request.client.host if request.client else "unknown"
                )
                logger.info(
                    f"[BANDWIDTH] {client_ip} | {request.method} {request.url.path} "
                    f"→ {response.status_code} | {size_kb}KB | {elapsed_ms}ms"
                )

        response.body_iterator = logging_iterator()
        return response


app.add_middleware(BandwidthLoggerMiddleware)

TEMP_DIR = Path(tempfile.gettempdir()) / "pdf_tool_backend"
TEMP_DIR.mkdir(exist_ok=True)


# ─── Security Helpers ──────────────────────────────────────────

async def safe_read_upload(file: UploadFile) -> bytes:
    """
    Read uploaded file into bytes with a hard size cap.
    Raises HTTP 413 if the file exceeds MAX_UPLOAD_BYTES.
    Avoids reading the entire file when the Content-Length header already
    reveals the file is too large.
    """
    # Quick check: if Content-Length header is present and already too large, reject early
    content_length = file.size  # FastAPI sets this from Content-Length
    if content_length is not None and content_length > MAX_UPLOAD_BYTES:
        raise HTTPException(
            status_code=413,
            detail=f"File too large. Maximum allowed size is {MAX_UPLOAD_MB} MB."
        )

    # Read in chunks so we can enforce the limit even without a Content-Length header
    chunks: list[bytes] = []
    total = 0
    chunk_size = 1024 * 64  # 64 KB chunks
    while True:
        chunk = await file.read(chunk_size)
        if not chunk:
            break
        total += len(chunk)
        if total > MAX_UPLOAD_BYTES:
            raise HTTPException(
                status_code=413,
                detail=f"File too large. Maximum allowed size is {MAX_UPLOAD_MB} MB."
            )
        chunks.append(chunk)
    return b"".join(chunks)


def safe_filename(user_filename: str | None, fallback: str = "output") -> str:
    """
    Sanitize a user-supplied filename for use in Content-Disposition headers.
    - Strips path separators (../  /  \\)
    - Keeps only alphanumeric, dash, underscore, dot
    - Truncates to 200 characters
    - Falls back to 'output' if result is empty
    """
    if not user_filename:
        return fallback
    # Take only the basename — no directory traversal possible
    name = Path(user_filename).name
    # Strip non-safe characters
    name = re.sub(r"[^\w.\-]", "_", name)
    # Truncate
    name = name[:200]
    return name if name else fallback


def safe_error(e: Exception, endpoint: str = "") -> HTTPException:
    """
    Logs the real exception server-side (including stack details).
    Returns a generic HTTP 500 that does NOT expose internal state to the user.
    """
    logger.error(f"[ERROR] endpoint={endpoint} type={type(e).__name__} msg={str(e)[:500]}")
    return HTTPException(
        status_code=500,
        detail="Unable to process this file. Please try again. If the problem persists, try a different or smaller file."
    )


def check_pdf_page_count(path: Path, endpoint: str = "") -> int:
    """
    Open a PDF and check its page count against MAX_PDF_PAGES.
    Raises HTTP 400 if too many pages. Returns page count on success.
    """
    try:
        import fitz
        doc = fitz.open(str(path))
        pages = len(doc)
        doc.close()
        if pages > MAX_PDF_PAGES:
            raise HTTPException(
                status_code=400,
                detail=f"PDF has too many pages ({pages}). Maximum allowed is {MAX_PDF_PAGES} pages."
            )
        return pages
    except HTTPException:
        raise
    except Exception as e:
        logger.warning(f"[PAGE-COUNT] Could not check page count for {path.name}: {e}")
        return 0  # If we can't read it, let the downstream handler deal with it


# ─── General Helpers ───────────────────────────────────────────

def temp_path(suffix: str) -> Path:
    return TEMP_DIR / f"{uuid.uuid4().hex}{suffix}"


def cleanup(*paths: Path):
    for p in paths:
        try:
            if p.exists():
                shutil.rmtree(p) if p.is_dir() else p.unlink()
        except Exception:
            pass


def hex_to_rgb(hx: str):
    hx = str(hx).lstrip("#")
    if len(hx) == 6:
        return tuple(int(hx[i:i+2], 16) / 255.0 for i in (0, 2, 4))
    return (0, 0, 0)


def has_non_latin(text: str) -> bool:
    return any(ord(c) > 0x00FF for c in text)


# ─── Root & Health ─────────────────────────────────────────────

@app.get("/")
def read_root():
    # Minimal root response — avoids large JSON payload draining bandwidth
    return {"status": "ok", "version": "2.1.0"}


# Ultra-lightweight health check — supports both HEAD (Uptime Robot) and GET.
# HEAD request: server sends ONLY headers (zero body bytes transferred).
# GET request: returns a tiny fixed-size JSON body.
# Rate-limited to 30 requests/minute per IP to block bots.
@app.api_route("/health", methods=["GET", "HEAD"])
@limiter.limit("30/minute")
def health(request: Request, response: Response):
    if request.method == "HEAD":
        # Zero bytes sent — perfect for Uptime Robot pings
        response.status_code = 200
        return Response(status_code=200)
    # GET: return a tiny JSON (only ~15 bytes after GZip)
    return JSONResponse(content={"status": "ok"}, status_code=200)


# ─── AUTO-OCR HELPER ────────────────────────────────────────────
# Smart scan detection: checks text density per page, not just total chars.
# If >= 50% pages are mostly blank of text → treat as scanned.

def is_pdf_scanned(inp_path: Path) -> bool:
    """Returns True if the PDF is a scanned/image-based document."""
    try:
        import fitz
        doc = fitz.open(inp_path)
        scanned_pages = 0
        total_pages = len(doc)
        if total_pages == 0:
            doc.close()
            return False
        for page in doc:
            text = page.get_text("text").strip()
            # A normal page with real text has at least 20 chars
            if len(text) < 20:
                scanned_pages += 1
        doc.close()
        # Consider scanned if >50% pages have very little text
        return scanned_pages / total_pages > 0.5
    except Exception:
        return False


async def run_hybrid_ocr(inp_path: Path, out_path: Path, lang: str = "eng+hin", deskew: bool = True, rotate: bool = True, force_ocr: bool = False):
    """
    Hybrid OCR: tries iLovePDF API first (if keys present), falls back to local OCRmyPDF.
    iLovePDF API key is ALWAYS read from server-side environment variables — never from frontend.
    """
    from fastapi.concurrency import run_in_threadpool

    ilovepdf_pub = os.environ.get("ILOVEPDF_PUBLIC_KEY", "").strip()
    ilovepdf_sec = os.environ.get("ILOVEPDF_SECRET_KEY", "").strip()

    if ilovepdf_pub and ilovepdf_sec:
        logger.info(f"[OCR] Using iLovePDF API for: {inp_path.name}")
        def run_ilovepdf():
            import concurrent.futures

            def _do_ilovepdf():
                from ilovepdf import PdfOcrTask
                task = PdfOcrTask(public_key=ilovepdf_pub, secret_key=ilovepdf_sec)
                task.add_file(str(inp_path))
                task.execute()

                out_dir = Path(tempfile.gettempdir()) / uuid.uuid4().hex
                out_dir.mkdir(exist_ok=True)
                task.download(str(out_dir))

                downloaded_files = list(out_dir.iterdir())
                if not downloaded_files:
                    shutil.rmtree(out_dir, ignore_errors=True)
                    raise RuntimeError("iLovePDF returned no output files")
                shutil.move(str(downloaded_files[0]), str(out_path))
                shutil.rmtree(out_dir, ignore_errors=True)

                # Guard: ensure the downloaded file is non-empty
                if not out_path.exists() or out_path.stat().st_size == 0:
                    raise RuntimeError("iLovePDF produced an empty output file")

            # Run with a timeout to prevent hung iLovePDF calls from blocking workers
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as executor:
                future = executor.submit(_do_ilovepdf)
                try:
                    future.result(timeout=90)  # 90-second timeout for iLovePDF
                except concurrent.futures.TimeoutError:
                    raise RuntimeError("iLovePDF OCR timed out after 90 seconds")

        try:
            await run_in_threadpool(run_ilovepdf)
            logger.info(f"[OCR] iLovePDF succeeded for: {inp_path.name}")
            return True
        except Exception as e:
            # Log the error without exposing the API key or internal state
            logger.error(f"[OCR] iLovePDF API failed: {type(e).__name__}. Falling back to local OCR.")
            # Clean up any partial output before local fallback
            if out_path.exists():
                out_path.unlink(missing_ok=True)

    logger.info(f"[OCR] Using local OCRmyPDF for: {inp_path.name}")
    try:
        import ocrmypdf
    except ImportError:
        raise RuntimeError("ocrmypdf is not installed. Run: pip install ocrmypdf")

    def run_ocrmypdf():
        ocrmypdf.ocr(
            str(inp_path), str(out_path), language=lang,
            deskew=deskew, rotate_pages=rotate,
            skip_text=not force_ocr, force_ocr=force_ocr,
            output_type="pdf",
            progress_bar=False
        )
    await run_in_threadpool(run_ocrmypdf)

    # Guard: local OCR must also produce a non-empty file
    if not out_path.exists() or out_path.stat().st_size == 0:
        raise RuntimeError("Local OCRmyPDF produced no output")

    logger.info(f"[OCR] Local OCRmyPDF succeeded for: {inp_path.name}")
    return True


def validate_ocr_output(original_pdf: Path, ocr_pdf: Path) -> tuple[bool, str]:
    """
    Validate the OCR output PDF against the original.
    Returns (True, "") on success, or (False, reason) on failure.
    Checks:
      1. File exists
      2. File size > 0
      3. Opens without error
      4. Page count matches original
      5. At least one page has a meaningful text layer
      6. No page renders as a completely blank (all-zero) pixmap
    """
    try:
        import fitz

        # 1. Existence
        if not ocr_pdf.exists():
            return False, "OCR output file does not exist"

        # 2. Size
        if ocr_pdf.stat().st_size == 0:
            return False, "OCR output file is empty (0 bytes)"

        # 3 & 4. Openability + page count
        try:
            ocr_doc = fitz.open(str(ocr_pdf))
        except Exception as open_err:
            return False, f"OCR output cannot be opened: {type(open_err).__name__}"

        try:
            orig_doc = fitz.open(str(original_pdf))
            orig_pages = len(orig_doc)
            orig_doc.close()
        except Exception:
            orig_pages = None

        ocr_pages = len(ocr_doc)
        if ocr_pages == 0:
            ocr_doc.close()
            return False, "OCR output PDF has 0 pages"

        if orig_pages is not None and ocr_pages != orig_pages:
            ocr_doc.close()
            return False, (
                f"OCR output page count ({ocr_pages}) does not match "
                f"original ({orig_pages})"
            )

        # 5. Text layer check — at least some text must be present across all pages
        total_text_len = 0
        blank_page_count = 0
        for page in ocr_doc:
            page_text = page.get_text("text").strip()
            total_text_len += len(page_text)

            # 6. Pixmap non-blank check (render at low DPI to save memory)
            try:
                pix = page.get_pixmap(dpi=36)
                if pix.width > 0 and pix.height > 0:
                    # Check if all pixels are white/near-white (sum of non-255 bytes)
                    import struct
                    samples = pix.samples
                    # samples is a bytes object — check if all bytes are 255 (white)
                    non_white = sum(1 for b in samples if b < 250)
                    if non_white == 0:
                        blank_page_count += 1
            except Exception:
                pass  # Pixmap failure is not fatal for validation

        ocr_doc.close()

        if total_text_len < 10:
            return False, (
                f"OCR output has no meaningful text layer "
                f"(only {total_text_len} chars extracted across all pages)"
            )

        if blank_page_count == ocr_pages:
            return False, "All pages in OCR output render as blank"

        return True, ""

    except Exception as e:
        return False, f"Validation error: {type(e).__name__}: {str(e)[:200]}"


async def ensure_auto_ocr(inp_path: Path, force: bool = False) -> Path:
    """
    If the PDF is scanned (or force=True), run hybrid OCR and return the path
    to a validated OCR PDF.

    SAFE ARCHITECTURE:
    - OCR is written to a TEMPORARY file (never overwrites the original until validated)
    - OCR output is fully validated before the original is replaced
    - If OCR fails or validation fails, raises HTTP 422 with a user-friendly message
    - NEVER silently returns the original scanned PDF to converters
    """
    if not force and not is_pdf_scanned(inp_path):
        logger.info(f"[Auto-OCR] PDF has text layer, OCR skipped: {inp_path.name}")
        return inp_path

    logger.info(f"[Auto-OCR] Scanned PDF detected, running hybrid OCR: {inp_path.name}")

    # Write OCR output to a NEW temp file — do NOT overwrite the original yet
    out_path = inp_path.with_name(inp_path.stem + "_ocr_candidate.pdf")
    if out_path.exists():
        out_path.unlink(missing_ok=True)

    ocr_succeeded = False
    ocr_error_detail = ""

    try:
        await run_hybrid_ocr(inp_path, out_path, force_ocr=force)
        ocr_succeeded = True
    except Exception as e:
        ocr_error_detail = f"{type(e).__name__}: {str(e)[:300]}"
        logger.error(f"[Auto-OCR] Both iLovePDF and local OCR failed: {ocr_error_detail}")

    if not ocr_succeeded:
        # Clean up any partial output
        out_path.unlink(missing_ok=True)
        raise HTTPException(
            status_code=422,
            detail=(
                "Unable to extract readable text from this scanned PDF. "
                "OCR processing failed. Please try a clearer scan or a different file."
            )
        )

    # Validate the OCR output BEFORE touching the original
    logger.info(f"[Auto-OCR] Validating OCR output: {out_path.name}")
    valid, reason = validate_ocr_output(inp_path, out_path)

    if not valid:
        logger.error(f"[Auto-OCR] OCR output validation failed: {reason}")
        out_path.unlink(missing_ok=True)
        raise HTTPException(
            status_code=422,
            detail=(
                "Unable to extract readable text from this scanned PDF. "
                "The OCR process completed but produced no usable text layer. "
                "Please try another PDF or a clearer scan."
            )
        )

    logger.info(f"[Auto-OCR] Validation passed. Replacing original with OCR output.")
    # Only NOW replace the original with the validated OCR output
    shutil.move(str(out_path), str(inp_path))
    return inp_path


def remove_scanned_page_backgrounds_from_docx(docx_path: str, pdf_path: str):
    """
    Safely removes ONLY the full-page scanned background images from a DOCX generated by pdf2docx,
    while preserving legitimate media (logos, photos, tables) using localized image fallbacks,
    and making the OCR text visible securely.
    """
    try:
        import zipfile
        import tempfile
        import os
        import shutil
        import copy
        import xml.etree.ElementTree as ET
        import fitz
        import cv2
        import numpy as np

        logger.info(f"[RemoveBackgrounds] Analyzing {pdf_path} for non-text regions...")
        
        # Phase 1: Detect non-text regions using PyMuPDF + OpenCV
        doc = fitz.open(pdf_path)
        regions_by_page = {}
        for page_num in range(len(doc)):
            page = doc[page_num]
            text_blocks = page.get_text("blocks")
            text_rects = []
            for b in text_blocks:
                if b[6] == 0:
                    r = fitz.Rect(b[:4])
                    r.x0 = max(0, r.x0 - 5)
                    r.y0 = max(0, r.y0 - 5)
                    r.x1 += 5
                    r.y1 += 5
                    text_rects.append(r)
            
            pix = page.get_pixmap(dpi=150)
            img = np.frombuffer(pix.samples, dtype=np.uint8).reshape(pix.h, pix.w, pix.n)
            if pix.n == 4:
                img = cv2.cvtColor(img, cv2.COLOR_RGBA2GRAY)
            elif pix.n == 3:
                img = cv2.cvtColor(img, cv2.COLOR_RGB2GRAY)
            else:
                img = np.copy(img)
                
            thresh = cv2.adaptiveThreshold(img, 255, cv2.ADAPTIVE_THRESH_GAUSSIAN_C, cv2.THRESH_BINARY_INV, 11, 2)
            kernel = cv2.getStructuringElement(cv2.MORPH_RECT, (15, 15))
            dilated = cv2.dilate(thresh, kernel, iterations=2)
            contours, _ = cv2.findContours(dilated, cv2.RETR_EXTERNAL, cv2.CHAIN_APPROX_SIMPLE)
            
            non_text_rects = []
            page_area = page.rect.get_area()
            
            for cnt in contours:
                x, y, w, h = cv2.boundingRect(cnt)
                x_pt = x * 72.0 / 150.0
                y_pt = y * 72.0 / 150.0
                w_pt = w * 72.0 / 150.0
                h_pt = h * 72.0 / 150.0
                
                if w_pt < 15 or h_pt < 15:
                    continue
                    
                rect = fitz.Rect(x_pt, y_pt, x_pt + w_pt, y_pt + h_pt)
                area = rect.get_area()
                if area == 0 or area > page_area * 0.90:
                    continue
                    
                overlap_area = 0
                for t_rect in text_rects:
                    intersect = rect.intersect(t_rect)
                    if not intersect.is_empty:
                        overlap_area += intersect.get_area()
                        
                if overlap_area / area > 0.4:
                    continue
                non_text_rects.append(rect)
                
            merged = []
            for r in non_text_rects:
                if not merged:
                    merged.append(r)
                    continue
                has_merged = False
                for i, m in enumerate(merged):
                    if r.intersects(m):
                        merged[i] = m.include_rect(r)
                        has_merged = True
                        break
                if not has_merged:
                    merged.append(r)
            regions_by_page[page_num] = merged
        doc.close()

        # Phase 2: Process the DOCX
        temp_dir = tempfile.mkdtemp()
        with zipfile.ZipFile(docx_path, 'r') as zip_ref:
            zip_ref.extractall(temp_dir)
            
        doc_xml_path = os.path.join(temp_dir, 'word', 'document.xml')
        if not os.path.exists(doc_xml_path):
            shutil.rmtree(temp_dir, ignore_errors=True)
            return

        namespaces = {'w': 'http://schemas.openxmlformats.org/wordprocessingml/2006/main',
                      'v': 'urn:schemas-microsoft-com:vml',
                      'wp': 'http://schemas.openxmlformats.org/drawingml/2006/wordprocessingDrawing',
                      'a': 'http://schemas.openxmlformats.org/drawingml/2006/main',
                      'pic': 'http://schemas.openxmlformats.org/drawingml/2006/picture',
                      'r': 'http://schemas.openxmlformats.org/officeDocument/2006/relationships',
                      'mc': 'http://schemas.openxmlformats.org/markup-compatibility/2006',
                      'w14': 'http://schemas.microsoft.com/office/word/2010/wordml',
                      'wp14': 'http://schemas.microsoft.com/office/word/2010/wordprocessingDrawing'}
                      
        for prefix, uri in namespaces.items():
            ET.register_namespace(prefix, uri)
            
        tree = ET.parse(doc_xml_path)
        root = tree.getroot()
        
        default_w, default_h = 11906, 16838
        for sect in root.findall('.//w:sectPr', namespaces):
            pgSz = sect.find('w:pgSz', namespaces)
            if pgSz is not None:
                w_attr = pgSz.get('{' + namespaces['w'] + '}w')
                h_attr = pgSz.get('{' + namespaces['w'] + '}h')
                if w_attr: default_w = int(w_attr)
                if h_attr: default_h = int(h_attr)
                
        W_pt = default_w / 20.0
        H_pt = default_h / 20.0
        page_area_pt = W_pt * H_pt
        
        changed = False
        removed_image_rIds = set()
        
        for parent in root.findall('.//w:drawing/..', namespaces):
            for drawing in parent.findall('w:drawing', namespaces):
                is_background = False
                coverage = 0.0
                
                extent = drawing.find('.//wp:extent', namespaces)
                if extent is not None:
                    cx = int(extent.get('cx', 0)) / 12700.0
                    cy = int(extent.get('cy', 0)) / 12700.0
                    img_area = cx * cy
                    coverage = img_area / page_area_pt
                    
                    if coverage >= 0.85:
                        blip = drawing.find('.//a:blip', namespaces)
                        if blip is not None:
                            is_background = True
                            
                if is_background:
                    logger.info(f"[RemoveBackgrounds] Background scan identified (coverage {coverage:.2f}). Applying fallbacks.")
                    blip = drawing.find('.//a:blip', namespaces)
                    if blip is not None:
                        embed_id = blip.get('{' + namespaces['r'] + '}embed')
                        if embed_id:
                            removed_image_rIds.add(embed_id)
                            
                    # Get regions for page 0 (assuming simple mapping for single background)
                    # For multi-page docs, pdf2docx usually outputs sequentially. We just use all regions.
                    # To be safer, we merge all non-text regions from all pages.
                    all_regions = []
                    for rs in regions_by_page.values():
                        all_regions.extend(rs)
                        
                    max_docpr_id = 1000
                    for d in root.findall('.//wp:docPr', namespaces):
                        try:
                            max_docpr_id = max(max_docpr_id, int(d.get('id', '0')))
                        except: pass
                        
                    for r in all_regions:
                        # PyMuPDF Rect normalizes to x0 <= x1, y0 <= y1
                        x_pt, y_pt, x1_pt, y1_pt = r.x0, r.y0, r.x1, r.y1
                        w_pt = x1_pt - x_pt
                        h_pt = y1_pt - y_pt
                        
                        fallback = copy.deepcopy(drawing)
                        
                        # UNIQUE ID fix for MS Word rendering bug
                        max_docpr_id += 1
                        docPr = fallback.find('.//wp:docPr', namespaces)
                        if docPr is not None:
                            docPr.set('id', str(max_docpr_id))
                            docPr.set('name', f'Fallback_{max_docpr_id}')
                            
                        # CRITICAL: Also update pic:cNvPr ID to prevent "unreadable content" corruption
                        cNvPr = fallback.find('.//pic:cNvPr', namespaces)
                        if cNvPr is not None:
                            cNvPr.set('id', str(max_docpr_id))
                            cNvPr.set('name', f'Fallback_Pic_{max_docpr_id}')
                        
                        # Adjust extent
                        ext = fallback.find('.//wp:extent', namespaces)
                        if ext is not None:
                            ext.set('cx', str(int(w_pt * 12700)))
                            ext.set('cy', str(int(h_pt * 12700)))
                        
                        spPr = fallback.find('.//pic:spPr', namespaces)
                        if spPr is not None:
                            xfrm = spPr.find('.//a:xfrm', namespaces)
                            if xfrm is not None:
                                ext2 = xfrm.find('.//a:ext', namespaces)
                                if ext2 is not None:
                                    ext2.set('cx', str(int(w_pt * 12700)))
                                    ext2.set('cy', str(int(h_pt * 12700)))
                        
                        # Adjust position
                        posH = fallback.find('.//wp:positionH/wp:posOffset', namespaces)
                        if posH is not None: posH.text = str(int(x_pt * 12700))
                        posV = fallback.find('.//wp:positionV/wp:posOffset', namespaces)
                        if posV is not None: posV.text = str(int(y_pt * 12700))
                        
                        # Add crop
                        blipFill = fallback.find('.//pic:blipFill', namespaces)
                        if blipFill is not None:
                            srcRect = ET.Element('{http://schemas.openxmlformats.org/drawingml/2006/main}srcRect')
                            srcRect.set('l', str(max(0, int(x_pt / W_pt * 100000))))
                            srcRect.set('t', str(max(0, int(y_pt / H_pt * 100000))))
                            srcRect.set('r', str(max(0, int((W_pt - x1_pt) / W_pt * 100000))))
                            srcRect.set('b', str(max(0, int((H_pt - y1_pt) / H_pt * 100000))))
                            blipFill.insert(1, srcRect)
                        
                        parent.append(fallback)
                    
                    parent.remove(drawing)
                    changed = True
                        
        if removed_image_rIds:
            for p in root.findall('.//w:p', namespaces):
                # 1. Remove paragraph level shading because pdf2docx uses it to mimic scan background
                pPr = p.find('w:pPr', namespaces)
                if pPr is not None:
                    p_shd = pPr.find('w:shd', namespaces)
                    if p_shd is not None:
                        pPr.remove(p_shd)
                        changed = True
                        
                # 2. Force text visibility unconditionally
                for r in p.findall('.//w:r', namespaces):
                    rPr = r.find('w:rPr', namespaces)
                    if rPr is not None:
                        vanish = rPr.find('w:vanish', namespaces)
                        color = rPr.find('w:color', namespaces)
                        shd = rPr.find('w:shd', namespaces)
                        
                        if vanish is not None:
                            rPr.remove(vanish)
                            changed = True
                            
                        # If color is white/hidden, force to auto
                        if color is not None:
                            val = color.get('{' + namespaces['w'] + '}val')
                            if val in ('FFFFFF', 'white', 'F2F2F2', 'f2f2f2'):
                                color.set('{' + namespaces['w'] + '}val', 'auto')
                                changed = True
                        else:
                            # If no color tag, sometimes it defaults to hidden if inherited. Ensure auto.
                            pass
                            
                        if shd is not None:
                            rPr.remove(shd)
                            changed = True

        if changed:
            tree.write(doc_xml_path, xml_declaration=True, encoding='UTF-8')
            
        with zipfile.ZipFile(docx_path, 'w', zipfile.ZIP_DEFLATED) as zip_out:
            for root_dir, dirs, files in os.walk(temp_dir):
                for file in files:
                    file_path = os.path.join(root_dir, file)
                    arcname = os.path.relpath(file_path, temp_dir)
                    zip_out.write(file_path, arcname)
                    
        shutil.rmtree(temp_dir, ignore_errors=True)
        try:
            shutil.copy(docx_path, "scratch/debug_after.docx")
        except Exception as e: pass
        logger.info(f"[RemoveBackgrounds] Successfully processed {docx_path}")
    except Exception as e:
        logger.error(f"[RemoveBackgrounds] Failed: {e}")

# ─── 1. PDF → Word ─────────────────────────────────────────────

@app.post("/convert/pdf-to-word")
@limiter.limit("10/minute")
async def pdf_to_word(request: Request, file: UploadFile = File(...), force_ocr: bool = Form(False)):
    try:
        from pdf2docx import Converter
    except ImportError:
        raise HTTPException(500, "Run: pip install pdf2docx PyMuPDF")

    inp = temp_path(".pdf")
    out = temp_path(".docx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "pdf-to-word")

        # 1. Run OCR if needed — raises HTTP 422 on failure, never silently continues
        was_scanned = force_ocr or is_pdf_scanned(inp)
        inp = await ensure_auto_ocr(inp, force=force_ocr)

        logger.info(f"[PDF2Word] Using pdf2docx for layout reconstruction: {inp.name}")
        cv = Converter(str(inp))
        cv.convert(str(out), multi_processing=False,
                   line_overlap_threshold=0.9,
                   min_svg_gap_dx=15.0)
        cv.close()

        if was_scanned:
            logger.info(f"[PDF2Word] Scanned/OCR PDF detected. Removing background images from DOCX.")
            remove_scanned_page_backgrounds_from_docx(str(out), inp)

        # 2. Validate DOCX output has meaningful content
        if not out.exists() or out.stat().st_size == 0:
            raise HTTPException(
                status_code=422,
                detail="Conversion produced an empty document. Please try a different file."
            )
        try:
            from docx import Document as _DocxDoc
            _check_doc = _DocxDoc(str(out))

            # pdf2docx stores text in XML shapes/textboxes (<w:t> tags), NOT in
            # doc.paragraphs. We must scan all XML text nodes to get a true count.
            _W_T = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}t"
            _total_chars = sum(
                len(elem.text or "")
                for elem in _check_doc.element.iter(_W_T)
            )

            # Fallback: also check raw file size — a non-trivial DOCX (>8KB above
            # the empty template baseline) almost certainly has real content.
            _docx_size_kb = out.stat().st_size / 1024
            _has_content_by_size = _docx_size_kb > 8

            logger.info(
                f"[PDF2Word] DOCX chars={_total_chars}, size={_docx_size_kb:.1f}KB"
            )

            if _total_chars < 5 and not _has_content_by_size:
                raise HTTPException(
                    status_code=422,
                    detail=(
                        "Conversion produced a document with no readable text. "
                        "If this is a scanned PDF, please ensure it is a clear scan."
                    )
                )
            logger.info(f"[PDF2Word] Output validated: {_total_chars} chars, {_docx_size_kb:.1f}KB")
        except HTTPException:
            raise
        except Exception as docx_val_err:
            logger.warning(f"[PDF2Word] DOCX validation skipped: {type(docx_val_err).__name__}")

        stem = safe_filename(file.filename, "document")
        out_name = Path(stem).stem + ".docx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "pdf-to-word")

@app.post("/convert/pdf-to-word-plaintext")
@limiter.limit("10/minute")
async def pdf_to_word_plaintext(request: Request, file: UploadFile = File(...), force_ocr: bool = Form(False)):
    try:
        import fitz
        from docx import Document
    except ImportError:
        raise HTTPException(500, "Run: pip install PyMuPDF python-docx")

    inp = temp_path(".pdf")
    out = temp_path(".docx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "pdf-to-word-plaintext")

        # 1. OCR if necessary — raises HTTP 422 on failure, never silently continues
        inp = await ensure_auto_ocr(inp, force=force_ocr)

        # 2. Extract plain text using fitz
        logger.info(f"[PDF2Word-Plaintext] Extracting text from {inp.name}")
        doc = fitz.open(str(inp))

        # 3. Write to a fresh Word file
        docx_doc = Document()
        total_text_added = 0
        for page in doc:
            text = page.get_text("text")
            if text.strip():
                docx_doc.add_paragraph(text.strip())
                docx_doc.add_page_break()
                total_text_added += len(text.strip())

        doc.close()

        # 4. Guard: do not return a blank document
        if total_text_added < 5:
            raise HTTPException(
                status_code=422,
                detail=(
                    "Unable to extract readable text from this PDF. "
                    "If this is a scanned document, please ensure OCR was performed successfully."
                )
            )

        docx_doc.save(str(out))

        stem = safe_filename(file.filename, "document")
        out_name = Path(stem).stem + "_plaintext.docx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "pdf-to-word-plaintext")



# ─── 2. PDF → Excel ────────────────────────────────────────────

@app.post("/convert/pdf-to-excel")
@limiter.limit("10/minute")
async def pdf_to_excel(request: Request, file: UploadFile = File(...), method: str = Form("auto"), force_ocr: bool = Form(False)):
    try:
        import pdfplumber, openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except ImportError:
        raise HTTPException(500, "Run: pip install pdfplumber openpyxl")

    inp = temp_path(".pdf")
    out = temp_path(".xlsx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "pdf-to-excel")
        # Raises HTTP 422 on OCR failure — never silently continues with image-only PDF
        inp = await ensure_auto_ocr(inp, force=force_ocr)

        wb = openpyxl.Workbook()
        wb.remove(wb.active)
        tables_found = False
        total_cells_written = 0

        # Try camelot (best for bordered tables)
        try:
            import camelot
            for flavor in (["lattice", "stream"] if method == "auto" else [method]):
                tables = camelot.read_pdf(str(inp), pages="all", flavor=flavor)
                if tables and len(tables) > 0:
                    tables_found = True
                    for i, tbl in enumerate(tables):
                        sheet_name = f"Table_{i+1}_p{tbl.page}"[:31]
                        ws = wb.create_sheet(sheet_name)
                        for r_idx, row in enumerate(tbl.df.values, 1):
                            for c_idx, val in enumerate(row, 1):
                                cell_val = str(val).strip() if val else ""
                                cell = ws.cell(r_idx, c_idx, cell_val)
                                if cell_val:
                                    total_cells_written += 1
                                if r_idx == 1:
                                    cell.font = Font(bold=True)
                                    cell.fill = PatternFill("solid", fgColor="D9E1F2")
                                cell.border = Border(
                                    left=Side(style="thin"), right=Side(style="thin"),
                                    top=Side(style="thin"), bottom=Side(style="thin")
                                )
                                cell.alignment = Alignment(wrap_text=True)
                    break
        except Exception as camelot_err:
            logger.warning(f"[PDF2Excel] Camelot failed, using pdfplumber fallback: {type(camelot_err).__name__}")

        # Fallback: pdfplumber
        if not tables_found:
            with pdfplumber.open(str(inp)) as pdf:
                for pg_num, pg in enumerate(pdf.pages, 1):
                    ws = wb.create_sheet(f"Page_{pg_num}")
                    row_idx = 1
                    # Use text-based geometry for borderless tables
                    pg_tables = pg.extract_tables(table_settings={
                        "vertical_strategy": "text",
                        "horizontal_strategy": "text",
                        "intersection_y_tolerance": 15
                    })
                    if pg_tables:
                        for tbl in pg_tables:
                            for row in tbl:
                                for c_idx, v in enumerate(row, 1):
                                    cell_val = str(v).strip() if v else ""
                                    ws.cell(row_idx, c_idx, cell_val)
                                    if cell_val:
                                        total_cells_written += 1
                                row_idx += 1
                            row_idx += 2  # gap between tables
                    else:
                        # Fallback for completely unstructured text: try to align by X-coordinates
                        words = pg.extract_words()
                        if words:
                            # Group words into lines based on Y coordinate
                            lines = {}
                            for w in words:
                                y = round(w['top'] / 5) * 5  # group within 5pts
                                lines.setdefault(y, []).append(w)

                            for y in sorted(lines.keys()):
                                line_words = sorted(lines[y], key=lambda w: w['x0'])
                                # simple column mapping: roughly every 50pts is a column
                                for w in line_words:
                                    col_idx = max(1, int(w['x0'] / 50) + 1)
                                    ws.cell(row_idx, col_idx, w['text'])
                                    total_cells_written += 1
                                row_idx += 1

        # Guard: do not return a blank workbook
        if total_cells_written == 0:
            raise HTTPException(
                status_code=422,
                detail=(
                    "No readable table or text data was detected in this PDF. "
                    "If this is a scanned document, ensure it contains legible text."
                )
            )

        logger.info(f"[PDF2Excel] Extracted {total_cells_written} non-empty cells")
        wb.save(str(out))

        stem = safe_filename(file.filename, "document")
        out_name = Path(stem).stem + ".xlsx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "pdf-to-excel")


# ─── 3. PDF → PPT ──────────────────────────────────────────────

@app.post("/convert/pdf-to-ppt")
@limiter.limit("10/minute")
async def pdf_to_ppt(request: Request, file: UploadFile = File(...), dpi: int = Form(150), force_ocr: bool = Form(False)):
    try:
        import fitz
        from pptx import Presentation
        from pptx.util import Inches, Pt
        from pptx.dml.color import RGBColor
        import io
    except ImportError:
        raise HTTPException(500, "Run: pip install PyMuPDF python-pptx")

    # Clamp DPI to safe range to prevent memory exhaustion
    dpi = max(MIN_DPI, min(dpi, MAX_DPI))

    inp = temp_path(".pdf")
    out = temp_path(".pptx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "pdf-to-ppt")
        # Raises HTTP 422 on OCR failure — never silently continues with image-only PDF
        inp = await ensure_auto_ocr(inp, force=force_ocr)

        doc = fitz.open(str(inp))
        prs = Presentation()
        slides_with_content = 0

        for page_num, page in enumerate(doc, 1):
            # Page dimension validation
            if page.rect.width <= 0 or page.rect.height <= 0:
                logger.warning(f"[PDF2PPT] Page {page_num} has zero dimensions, skipping")
                continue

            w_in = page.rect.width / 72
            h_in = page.rect.height / 72
            prs.slide_width = Inches(w_in)
            prs.slide_height = Inches(h_in)
            slide = prs.slides.add_slide(prs.slide_layouts[6])

            # ── STEP 1: Render the full page as a high-quality image (primary visual) ──
            # This is the PRIMARY source of visual content for scanned PDFs.
            # The rendered page image is ALWAYS inserted first regardless of OCR.
            pix = page.get_pixmap(dpi=dpi)

            # Validate the rendered pixmap is not blank
            pixmap_is_blank = False
            try:
                samples = pix.samples
                non_white = sum(1 for b in samples if b < 250)
                if non_white == 0:
                    pixmap_is_blank = True
                    logger.warning(f"[PDF2PPT] Page {page_num} rendered as a blank pixmap")
            except Exception as pix_check_err:
                logger.warning(f"[PDF2PPT] Pixmap check failed for page {page_num}: {pix_check_err}")

            if pixmap_is_blank:
                # Abort entire conversion — do not produce a deck of blank slides
                doc.close()
                raise HTTPException(
                    status_code=422,
                    detail=(
                        f"Page {page_num} of this PDF rendered as completely blank. "
                        "Conversion aborted to avoid producing a blank presentation. "
                        "Please try a different or clearer file."
                    )
                )

            # Insert the rendered page image covering the full slide
            img_bytes = pix.tobytes("png")
            slide.shapes.add_picture(
                io.BytesIO(img_bytes),
                Inches(0), Inches(0),
                Inches(w_in), Inches(h_in)
            )
            slides_with_content += 1

            # ── STEP 2: Overlay OCR text blocks as an invisible/searchable text layer ──
            # This is ADDITIVE and subordinate to the image — it provides
            # copy-paste / search functionality when OCR text is available.
            blocks = page.get_text("dict").get("blocks", [])
            for block in blocks:
                if block["type"] == 0:  # Text block from OCR
                    x0, y0, x1, y1 = block["bbox"]
                    block_w = max((x1 - x0) / 72, 0.3)
                    block_h = max((y1 - y0) / 72, 0.2)
                    try:
                        tb = slide.shapes.add_textbox(
                            Inches(x0 / 72), Inches(y0 / 72),
                            Inches(block_w), Inches(block_h)
                        )
                        tf = tb.text_frame
                        tf.word_wrap = True
                        for line_idx, line in enumerate(block.get("lines", [])):
                            if line_idx > 0:
                                p = tf.add_paragraph()
                            else:
                                p = tf.paragraphs[0]

                            for span in line.get("spans", []):
                                txt = span.get("text", "")
                                if not txt:
                                    continue
                                run = p.add_run()
                                run.text = txt
                                run.font.size = Pt(max(6, span.get("size", 12)))
                                # Make text invisible (transparent) over the image
                                # so only the image is visible, but text is searchable
                                run.font.color.rgb = RGBColor(255, 255, 255)
                                run.font.color.rgb = RGBColor(
                                    (span.get("color", 0) >> 16) & 255,
                                    (span.get("color", 0) >> 8) & 255,
                                    span.get("color", 0) & 255
                                )
                                flags = span.get("flags", 0)
                                if flags & 2 ** 4:
                                    run.font.bold = True
                                if flags & 2 ** 1:
                                    run.font.italic = True
                    except Exception as tb_err:
                        # Text overlay is best-effort — never block slide creation
                        logger.debug(f"[PDF2PPT] Text overlay failed for block on page {page_num}: {tb_err}")

        doc.close()

        # Guard: do not return a presentation with no visible content
        if slides_with_content == 0:
            raise HTTPException(
                status_code=422,
                detail="No pages could be rendered from this PDF. Conversion aborted."
            )

        logger.info(f"[PDF2PPT] Generated {slides_with_content} slides with page images")
        prs.save(str(out))

        stem = safe_filename(file.filename, "presentation")
        out_name = Path(stem).stem + ".pptx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "pdf-to-ppt")


# ─── 4. Office → PDF (LibreOffice) ────────────────────────────

def _find_libreoffice() -> str:
    """Find LibreOffice executable. Returns path or raises RuntimeError."""
    candidates = [
        "soffice", "libreoffice",
        "/Applications/LibreOffice.app/Contents/MacOS/soffice",
        "/usr/bin/soffice",
    ]
    for c in candidates:
        if shutil.which(c) or Path(c).exists():
            return c
    raise RuntimeError("LibreOffice not found.")


def libreoffice_convert(inp: Path, out_dir: Path) -> Path:
    """
    Convert a file to PDF using LibreOffice.
    Arguments are passed as a list (no shell=True) to prevent injection.
    """
    lo = _find_libreoffice()
    r = subprocess.run(
        [lo, "--headless", "--convert-to", "pdf", "--outdir", str(out_dir), str(inp)],
        capture_output=True, text=True, timeout=MAX_PROCESSING_SECONDS
    )
    if r.returncode != 0:
        # Log stderr internally but only expose a generic message
        logger.error(f"[LibreOffice] Conversion failed: {r.stderr[:300]}")
        raise RuntimeError("LibreOffice conversion failed.")
    out_pdf = out_dir / (inp.stem + ".pdf")
    if not out_pdf.exists():
        raise RuntimeError("LibreOffice output PDF not found.")
    return out_pdf


ALLOWED_OFFICE_EXTENSIONS = {".docx", ".doc", ".pptx", ".ppt", ".xlsx", ".xls", ".odt", ".odp"}


@app.post("/convert/office-to-pdf")
@limiter.limit("10/minute")
async def office_to_pdf(request: Request, file: UploadFile = File(...)):
    ext = Path(safe_filename(file.filename or "")).suffix.lower()
    if ext not in ALLOWED_OFFICE_EXTENSIONS:
        raise HTTPException(400, f"Unsupported format. Allowed: {', '.join(ALLOWED_OFFICE_EXTENSIONS)}")

    inp = temp_path(ext)
    out_dir = TEMP_DIR / uuid.uuid4().hex
    out_dir.mkdir(exist_ok=True)
    try:
        inp.write_bytes(await safe_read_upload(file))

        # Acquire semaphore to limit concurrent LibreOffice processes
        async with LIBREOFFICE_SEMAPHORE:
            from fastapi.concurrency import run_in_threadpool
            out_pdf = await run_in_threadpool(libreoffice_convert, inp, out_dir)

        stem = safe_filename(file.filename, "document")
        out_name = Path(stem).stem + ".pdf"
        return FileResponse(
            str(out_pdf), media_type="application/pdf",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out_dir)
        )
    except HTTPException:
        cleanup(inp, out_dir)
        raise
    except Exception as e:
        cleanup(inp, out_dir)
        raise safe_error(e, "office-to-pdf")


@app.post("/convert/make-searchable")
@limiter.limit("5/minute")
async def make_searchable(
    request: Request,
    file: UploadFile = File(...),
    lang: str = Form("hin+eng"),
    deskew: bool = Form(True),
    rotate: bool = Form(True),
    force_ocr: bool = Form(False)
):
    inp = temp_path(".pdf")
    out = temp_path("_searchable.pdf")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "make-searchable")

        # Smart Check: If it's not scanned and user didn't force OCR, skip to save credits!
        if not force_ocr and not is_pdf_scanned(inp):
            logger.info(f"[OCR] Skipping OCR — PDF already has text. (Credit Saver)")
            # Just return the original file
            stem = safe_filename(file.filename, "document")
            out_name = Path(stem).stem + "_searchable.pdf"
            return FileResponse(
                str(inp), media_type="application/pdf",
                filename=out_name,
                background=BackgroundTask(cleanup, inp)
            )

        await run_hybrid_ocr(inp, out, lang=lang, deskew=deskew, rotate=rotate, force_ocr=force_ocr)
        stem = safe_filename(file.filename, "document")
        out_name = Path(stem).stem + "_searchable.pdf"
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "make-searchable")


# ─── 6. PPT → Word ────────────────────────────────────────────

@app.post("/convert/ppt-to-word")
@limiter.limit("10/minute")
async def ppt_to_word(request: Request, file: UploadFile = File(...)):
    try:
        from pptx import Presentation as Pptx
        from docx import Document
        from docx.shared import Pt as DPt
    except ImportError:
        raise HTTPException(500, "Run: pip install python-pptx python-docx")

    inp = temp_path(".pptx")
    out = temp_path(".docx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        prs = Pptx(str(inp))
        doc = Document()
        doc.add_heading("Presentation Content", 0)

        for i, slide in enumerate(prs.slides, 1):
            doc.add_heading(f"Slide {i}", 1)
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        txt = para.text.strip()
                        if txt:
                            p = doc.add_paragraph(txt)
                            if para.runs and para.runs[0].font.size:
                                p.runs[0].font.size = para.runs[0].font.size

        doc.save(str(out))
        stem = safe_filename(file.filename, "document")
        out_name = Path(stem).stem + ".docx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "ppt-to-word")


# ─── 7. Word → PPT (improved: headings→slides, bullets, tables) ──

@app.post("/convert/word-to-ppt")
@limiter.limit("10/minute")
async def word_to_ppt(request: Request, file: UploadFile = File(...)):
    try:
        from docx import Document
        from docx.oxml.ns import qn
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
        from pptx.enum.text import PP_ALIGN
    except ImportError:
        raise HTTPException(500, "Run: pip install python-docx python-pptx")

    inp = temp_path(".docx")
    out = temp_path(".pptx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        doc = Document(str(inp))
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        TITLE_LAYOUT   = prs.slide_layouts[0]   # Title Slide
        CONTENT_LAYOUT = prs.slide_layouts[1]   # Title + Content
        BLANK_LAYOUT   = prs.slide_layouts[6]   # Blank

        def _add_content_slide(title_text: str) -> any:
            sl = prs.slides.add_slide(CONTENT_LAYOUT)
            sl.shapes.title.text = title_text[:200]
            return sl

        def _tf_for_slide(sl) -> any:
            for ph in sl.placeholders:
                if ph.placeholder_format.idx == 1:
                    return ph.text_frame
            return None

        def _hex_to_rgb(hex_str: str):
            h = hex_str.lstrip('#')
            if len(h) == 6:
                return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))
            return None

        current_slide: any = None
        current_tf: any = None
        first_slide = True

        for elem in doc.element.body:
            tag = elem.tag.split('}')[-1] if '}' in elem.tag else elem.tag

            if tag == 'p':
                from docx import Document as _D
                from docx.text.paragraph import Paragraph as _P
                para = _P(elem, doc)
                txt  = para.text.strip()
                style = para.style.name if para.style else ''

                is_h1 = 'Heading 1' in style or 'Title' in style
                is_h2 = 'Heading 2' in style
                is_h3 = 'Heading 3' in style or 'Heading 4' in style

                if is_h1:
                    current_slide = _add_content_slide(txt or 'Slide')
                    current_tf    = _tf_for_slide(current_slide)
                    if current_tf:
                        current_tf.clear()
                    first_slide = False
                    continue

                if first_slide and txt:
                    current_slide = _add_content_slide(txt or 'Content')
                    current_tf    = _tf_for_slide(current_slide)
                    if current_tf:
                        current_tf.clear()
                    first_slide = False
                    continue

                if not txt:
                    continue

                if current_slide is None:
                    current_slide = _add_content_slide('Content')
                    current_tf    = _tf_for_slide(current_slide)
                    if current_tf:
                        current_tf.clear()

                if current_tf:
                    # determine indent level
                    level = 0
                    if is_h2:     level = 0
                    elif is_h3:   level = 1
                    elif para.paragraph_format.left_indent and para.paragraph_format.left_indent > 0:
                        level = min(4, int(para.paragraph_format.left_indent / 360000))

                    # add paragraph with per-run formatting
                    pptx_para = current_tf.add_paragraph()
                    pptx_para.level = level
                    added_run = False
                    for run in para.runs:
                        run_text = run.text
                        if not run_text:
                            continue
                        pptx_run = pptx_para.add_run()
                        pptx_run.text = run_text
                        pptx_run.font.bold   = run.bold
                        pptx_run.font.italic = run.italic
                        if run.font.size:
                            pptx_run.font.size = run.font.size
                        else:
                            pptx_run.font.size = Pt(20 if is_h2 else 18)
                        if run.font.color and run.font.color.type:
                            try:
                                col = run.font.color.rgb
                                pptx_run.font.color.rgb = RGBColor(col.red, col.green, col.blue)
                            except Exception:
                                pass
                        added_run = True
                    if not added_run:
                        pptx_run = pptx_para.add_run()
                        pptx_run.text = txt
                        pptx_run.font.size = Pt(20 if is_h2 else 18)

            elif tag == 'tbl':
                # Word table → insert as text table on a new slide
                from docx.table import Table as _T
                tbl = _T(elem, doc)
                if current_slide is None:
                    current_slide = _add_content_slide('Table')
                    current_tf    = _tf_for_slide(current_slide)
                    if current_tf:
                        current_tf.clear()
                # Render table rows as bullet list (PPTX doesn't have native editable tables in placeholders easily)
                if current_tf:
                    for ri, row in enumerate(tbl.rows):
                        cells_txt = ' | '.join(c.text.strip() for c in row.cells if c.text.strip())
                        if not cells_txt:
                            continue
                        pp = current_tf.add_paragraph()
                        pp.level = 1
                        r = pp.add_run()
                        r.text = cells_txt
                        r.font.bold = (ri == 0)
                        r.font.size = Pt(14)

        if len(prs.slides) == 0:
            prs.slides.add_slide(BLANK_LAYOUT)
        prs.save(str(out))

        stem     = safe_filename(file.filename, "presentation")
        out_name = Path(stem).stem + ".pptx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "word-to-ppt")


# ─── 8. Excel → Word ────────────────────────────────────────────

@app.post("/convert/excel-to-word")
@limiter.limit("10/minute")
async def excel_to_word(request: Request, file: UploadFile = File(...)):
    try:
        import openpyxl
        from docx import Document
        from docx.shared import Pt, RGBColor, Inches
        from docx.enum.text import WD_ALIGN_PARAGRAPH
        from docx.oxml.ns import qn
        from docx.oxml import OxmlElement
    except ImportError:
        raise HTTPException(500, "Run: pip install openpyxl python-docx")

    inp = temp_path(".xlsx")
    out = temp_path(".docx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        wb = openpyxl.load_workbook(str(inp), data_only=True)
        doc = Document()

        # Document margins
        for section in doc.sections:
            section.top_margin    = Inches(0.75)
            section.bottom_margin = Inches(0.75)
            section.left_margin   = Inches(0.75)
            section.right_margin  = Inches(0.75)

        for sheet_idx, sheet_name in enumerate(wb.sheetnames):
            ws = wb[sheet_name]
            if sheet_idx > 0:
                doc.add_page_break()

            h = doc.add_heading(sheet_name, level=1)
            h.runs[0].font.color.rgb = RGBColor(0x1a, 0x3a, 0x8f)

            # Determine data range (skip completely empty rows/cols)
            rows_data = []
            for row in ws.iter_rows(values_only=True):
                if any(c is not None and str(c).strip() for c in row):
                    rows_data.append([str(c) if c is not None else '' for c in row])
            if not rows_data:
                doc.add_paragraph('(Empty sheet)')
                continue

            # Trim trailing empty columns
            max_cols = max((len([c for c in r if c.strip()]) for r in rows_data), default=0)
            max_cols = max(max_cols, 1)
            rows_data = [r[:max_cols] for r in rows_data]
            # Pad short rows
            rows_data = [r + [''] * (max_cols - len(r)) for r in rows_data]

            table = doc.add_table(rows=len(rows_data), cols=max_cols)
            table.style = 'Table Grid'

            for ri, row in enumerate(rows_data):
                is_header = ri == 0
                tr = table.rows[ri]
                for ci, val in enumerate(row):
                    cell = tr.cells[ci]
                    cell.text = val
                    run = cell.paragraphs[0].runs[0] if cell.paragraphs[0].runs else cell.paragraphs[0].add_run(val)
                    run.bold = is_header
                    run.font.size = Pt(10 if not is_header else 11)
                    # Header shading
                    if is_header:
                        from docx.oxml.ns import qn
                        from docx.oxml import OxmlElement
                        tc = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        shd = OxmlElement('w:shd')
                        shd.set(qn('w:val'), 'clear')
                        shd.set(qn('w:color'), 'auto')
                        shd.set(qn('w:fill'), '1A3A8F')
                        tcPr.append(shd)
                        run.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)
                    elif ri % 2 == 0:
                        tc = cell._tc
                        tcPr = tc.get_or_add_tcPr()
                        shd = OxmlElement('w:shd')
                        shd.set(qn('w:val'), 'clear')
                        shd.set(qn('w:color'), 'auto')
                        shd.set(qn('w:fill'), 'EFF6FF')
                        tcPr.append(shd)

            doc.add_paragraph()

        doc.save(str(out))
        stem     = safe_filename(file.filename, "document")
        out_name = Path(stem).stem + ".docx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.wordprocessingml.document",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "excel-to-word")


# ─── 9. Excel → PPT ────────────────────────────────────────────

@app.post("/convert/excel-to-ppt")
@limiter.limit("10/minute")
async def excel_to_ppt(request: Request, file: UploadFile = File(...)):
    try:
        import openpyxl
        from pptx import Presentation
        from pptx.util import Inches, Pt, Emu
        from pptx.dml.color import RGBColor
    except ImportError:
        raise HTTPException(500, "Run: pip install openpyxl python-pptx")

    inp = temp_path(".xlsx")
    out = temp_path(".pptx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        wb = openpyxl.load_workbook(str(inp), data_only=True)
        prs = Presentation()
        prs.slide_width  = Inches(13.33)
        prs.slide_height = Inches(7.5)

        TITLE_LAYOUT   = prs.slide_layouts[0]
        CONTENT_LAYOUT = prs.slide_layouts[5]  # blank layout to draw manually
        TABLE_LAYOUT   = prs.slide_layouts[6]  # blank

        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]

            # Collect non-empty rows
            rows_data = []
            for row in ws.iter_rows(values_only=True):
                row_vals = [str(c) if c is not None else '' for c in row]
                if any(v.strip() for v in row_vals):
                    rows_data.append(row_vals)
            if not rows_data:
                continue

            # Trim trailing empty columns
            max_cols = max((len([c for c in r if c.strip()]) for r in rows_data), default=1)
            rows_data = [r[:max_cols] for r in rows_data]
            rows_data = [r + [''] * (max_cols - len(r)) for r in rows_data]

            # One slide per sheet — add title + table shape
            slide = prs.slides.add_slide(TABLE_LAYOUT)

            # Title text box at top
            title_box = slide.shapes.add_textbox(
                Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.65)
            )
            tf = title_box.text_frame
            tf.word_wrap = False
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = sheet_name
            run.font.size = Pt(24)
            run.font.bold = True
            run.font.color.rgb = RGBColor(0x1a, 0x3a, 0x8f)

            # Add pptx table
            max_rows = min(len(rows_data), 30)  # cap at 30 rows per slide
            table_rows = rows_data[:max_rows]
            n_rows = len(table_rows)
            n_cols = max_cols

            left   = Inches(0.3)
            top    = Inches(0.9)
            width  = Inches(12.7)
            height = Inches(6.4)

            tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
            tbl = tbl_shape.table

            for ri, row in enumerate(table_rows):
                is_header = ri == 0
                for ci, val in enumerate(row):
                    cell = tbl.cell(ri, ci)
                    cell.text = val[:100]
                    tf2 = cell.text_frame
                    tf2.word_wrap = True
                    run2 = tf2.paragraphs[0].runs[0] if tf2.paragraphs[0].runs else tf2.paragraphs[0].add_run()
                    run2.text = val[:100]
                    run2.font.size = Pt(11 if is_header else 10)
                    run2.font.bold = is_header
                    # Header background
                    from pptx.oxml.ns import qn as pqn
                    from lxml import etree
                    tc = cell._tc
                    tcPr = tc.get_or_add_tcPr()
                    solidFill = etree.SubElement(tcPr, pqn('a:solidFill'))
                    srgbClr = etree.SubElement(solidFill, pqn('a:srgbClr'))
                    srgbClr.set('val', '1A3A8F' if is_header else ('EFF6FF' if ri % 2 == 0 else 'FFFFFF'))
                    if is_header:
                        run2.font.color.rgb = RGBColor(0xFF, 0xFF, 0xFF)

            # If sheet had more rows, add continuation slide
            if len(rows_data) > max_rows:
                remaining = rows_data[max_rows:]
                slide2 = prs.slides.add_slide(TABLE_LAYOUT)
                tb2 = slide2.shapes.add_textbox(Inches(0.3), Inches(0.15), Inches(12.5), Inches(0.65))
                tf3 = tb2.text_frame
                r3 = tf3.paragraphs[0].add_run()
                r3.text = f"{sheet_name} (continued)"
                r3.font.size = Pt(22); r3.font.bold = True
                r3.font.color.rgb = RGBColor(0x1a, 0x3a, 0x8f)

                n2 = min(len(remaining), 30)
                tbl2_shape = slide2.shapes.add_table(n2, n_cols, left, top, width, height)
                tbl2 = tbl2_shape.table
                for ri, row in enumerate(remaining[:n2]):
                    for ci, val in enumerate(row):
                        tbl2.cell(ri, ci).text = val[:100]

        if len(prs.slides) == 0:
            prs.slides.add_slide(prs.slide_layouts[6])

        prs.save(str(out))
        stem     = safe_filename(file.filename, "presentation")
        out_name = Path(stem).stem + ".pptx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "excel-to-ppt")


# ─── 10. PPT → Excel ────────────────────────────────────────────

@app.post("/convert/ppt-to-excel")
@limiter.limit("10/minute")
async def ppt_to_excel(request: Request, file: UploadFile = File(...)):
    try:
        from pptx import Presentation as Pptx
        import openpyxl
        from openpyxl.styles import Font, PatternFill, Alignment, Border, Side
    except ImportError:
        raise HTTPException(500, "Run: pip install python-pptx openpyxl")

    inp = temp_path(".pptx")
    out = temp_path(".xlsx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        prs = Pptx(str(inp))
        wb  = openpyxl.Workbook()
        wb.remove(wb.active)

        HDR_FILL  = PatternFill("solid", fgColor="1A3A8F")
        ALT_FILL  = PatternFill("solid", fgColor="EFF6FF")
        HDR_FONT  = Font(bold=True, color="FFFFFF", size=11)
        BODY_FONT = Font(size=10)
        THIN      = Side(style="thin")
        BORDER    = Border(left=THIN, right=THIN, top=THIN, bottom=THIN)
        WRAP      = Alignment(wrap_text=True, vertical="top")

        for slide_idx, slide in enumerate(prs.slides, 1):
            ws = wb.create_sheet(f"Slide_{slide_idx}")
            row_idx = 1

            # Slide title in first row
            title_text = ""
            for shape in slide.shapes:
                if shape.has_text_frame and shape.shape_type == 13:
                    title_text = shape.text_frame.text.strip()
                    break
            if not title_text:
                for shape in slide.shapes:
                    if hasattr(shape, 'placeholder_format') and shape.placeholder_format:
                        if shape.placeholder_format.idx == 0:
                            title_text = shape.text_frame.text.strip() if shape.has_text_frame else ''
                            break

            if title_text:
                cell = ws.cell(row_idx, 1, title_text)
                cell.font  = Font(bold=True, size=14, color="1A3A8F")
                cell.fill  = PatternFill("solid", fgColor="D1E4FF")
                cell.alignment = WRAP
                ws.merge_cells(start_row=row_idx, start_column=1, end_row=row_idx, end_column=4)
                row_idx += 2

            # Extract tables from slide
            for shape in slide.shapes:
                if shape.has_table:
                    tbl = shape.table
                    for ri, row in enumerate(tbl.rows):
                        is_header = ri == 0
                        for ci, cell in enumerate(row.cells):
                            c = ws.cell(row_idx, ci + 1, cell.text.strip())
                            c.font      = HDR_FONT if is_header else BODY_FONT
                            c.fill      = HDR_FILL if is_header else (ALT_FILL if row_idx % 2 == 0 else PatternFill())
                            c.border    = BORDER
                            c.alignment = WRAP
                        row_idx += 1
                    row_idx += 1  # gap after table

            # Extract text frames (non-title)
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                is_title_shape = (hasattr(shape, 'placeholder_format') and
                                  shape.placeholder_format and
                                  shape.placeholder_format.idx == 0)
                if is_title_shape:
                    continue
                for para in shape.text_frame.paragraphs:
                    txt = para.text.strip()
                    if txt:
                        c = ws.cell(row_idx, 1, txt)
                        c.font      = BODY_FONT
                        c.alignment = WRAP
                        row_idx += 1

            # Auto-column widths
            for col in ws.columns:
                max_len = 0
                col_letter = col[0].column_letter
                for cell in col:
                    try:
                        if cell.value:
                            max_len = max(max_len, len(str(cell.value)))
                    except Exception:
                        pass
                ws.column_dimensions[col_letter].width = min(max(max_len + 2, 12), 50)

        if not wb.sheetnames:
            wb.create_sheet("Sheet1")

        wb.save(str(out))
        stem     = safe_filename(file.filename, "spreadsheet")
        out_name = Path(stem).stem + ".xlsx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "ppt-to-excel")


# ─── 11. Edit PDF (PyMuPDF redact + rewrite) ────────────────────


@app.post("/edit-pdf")
@limiter.limit("15/minute")
async def edit_pdf_endpoint(request: Request, file: UploadFile = File(...), edits: str = Form(...)):
    import fitz
    inp = temp_path(".pdf")
    out = temp_path(".pdf")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "edit-pdf")

        # Validate edits JSON — cap payload size to prevent memory abuse
        if len(edits) > 512_000:  # 512 KB max for edits JSON
            raise HTTPException(400, "Edit data is too large.")
        edit_actions = json.loads(edits)
        if not isinstance(edit_actions, list):
            raise HTTPException(400, "Invalid edit format.")

        doc = fitz.open(str(inp))

        # ── Font path resolver for Devanagari/Hindi ──────────────────────────
        _fonts_dir = Path(__file__).parent.parent / "public" / "fonts"

        def _resolve_font_path(bold: bool) -> Path | None:
            """Return the correct NotoSansDevanagari font path, or None if not found."""
            key = "NotoSansDevanagari-Bold.ttf" if bold else "NotoSansDevanagari-Regular.ttf"
            candidate = _fonts_dir / key
            if candidate.exists():
                return candidate
            candidate2 = Path(__file__).parent / "fonts" / key
            if candidate2.exists():
                return candidate2
            return None

        _registered_fonts: dict[int, set] = {}

        for action in edit_actions:
            page_idx = action.get("page", 0)
            if page_idx >= len(doc):
                continue
            page = doc[page_idx]

            # PyMuPDF uses TOP-LEFT origin (same as browser canvas).
            # Frontend divides canvas coords by viewport.scale before sending,
            # so x,y,w,h are already in PDF points with y measured from top.
            x   = float(action["x"])
            y   = float(action["y"])   # top of text bounding box, from top of page
            w   = float(action["w"])
            h   = float(action["h"])
            new_text = action.get("newText", "")

            # Redact: add small padding around box to fully cover existing text
            rect = fitz.Rect(x - 2, y - 2, x + w + 2, y + h + 2)
            
            # Smart Background Reconstruction
            has_native_text = False
            for b in page.get_text("blocks", clip=rect):
                if b[6] == 0 and b[4].strip():
                    has_native_text = True
                    break
                    
            if has_native_text:
                # Native PDF: remove text objects, preserve images underneath
                page.add_redact_annot(rect, cross_out=False)
                page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_NONE)
            else:
                # Scanned PDF: sample edge pixels to get dominant background color
                try:
                    pm = page.get_pixmap(clip=rect)
                    from collections import Counter
                    edges = []
                    for px in range(pm.width):
                        edges.append(pm.pixel(px, 0))
                        edges.append(pm.pixel(px, pm.height - 1))
                    for py in range(pm.height):
                        edges.append(pm.pixel(0, py))
                        edges.append(pm.pixel(pm.width - 1, py))
                    if edges:
                        mc = Counter(edges).most_common(1)[0][0]
                        # Handling RGB or RGBA tuples
                        bg_color = (mc[0]/255.0, mc[1]/255.0, mc[2]/255.0)
                    else:
                        bg_color = (1, 1, 1)
                except Exception:
                    bg_color = (1, 1, 1)
                
                page.add_redact_annot(rect, fill=bg_color)
                page.apply_redactions(images=fitz.PDF_REDACT_IMAGE_PIXELS)

            if not new_text:
                continue

            fg_color  = hex_to_rgb(action.get("color", "#000000"))
            orig_size = float(action.get("fontSize", 12))
            if orig_size < 4:
                orig_size = 12
            is_bold   = bool(action.get("bold"))
            is_italic = bool(action.get("italic"))

            # Latin font selection
            if is_bold and is_italic:
                fontname = "hebi"
            elif is_bold:
                fontname = "hebo"
            elif is_italic:
                fontname = "heit"
            else:
                fontname = "helv"

            # ── Hindi/Devanagari font handling ────────────────────────────────
            if has_non_latin(new_text):
                deva_path = _resolve_font_path(is_bold)
                if deva_path is not None:
                    fontname = "NotoDevaB" if is_bold else "NotoDeva"
                    page_fonts = _registered_fonts.setdefault(page_idx, set())
                    if fontname not in page_fonts:
                        try:
                            font_bytes = deva_path.read_bytes()
                            page.insert_font(fontname=fontname, fontbuffer=font_bytes)
                            page_fonts.add(fontname)
                            logger.info(f"[EditPDF] Registered Devanagari font '{fontname}' on page {page_idx}")
                        except Exception as fe:
                            logger.warning(f"[EditPDF] Devanagari font register failed: {fe}. Using helv.")
                            fontname = "hebo" if is_bold else "helv"
                else:
                    logger.warning(f"[EditPDF] Devanagari font not found at {_fonts_dir}.")
            # ─────────────────────────────────────────────────────────────────

            # Use insert_textbox for multiline wrapping and proper baseline handling
            text_rect = fitz.Rect(x, y, x + w, y + h)
            align = fitz.TEXT_ALIGN_LEFT
            
            rc = -1
            final_size = orig_size
            
            # Try to fit the text box, reducing font size if it doesn't fit
            while final_size >= 4:
                rc = page.insert_textbox(
                    text_rect,
                    new_text,
                    fontsize=final_size,
                    fontname=fontname,
                    color=fg_color,
                    align=align
                )
                if rc >= 0:
                    break
                final_size -= 0.5
                
            # If it still didn't fit, force it at the minimum size
            if rc < 0:
                page.insert_textbox(
                    text_rect,
                    new_text,
                    fontsize=final_size,
                    fontname=fontname,
                    color=fg_color,
                    align=align
                )

        doc.save(str(out), garbage=3, deflate=True)
        doc.close()

        stem = safe_filename(file.filename, "document")
        out_name = "edited_" + stem
        return FileResponse(str(out), media_type="application/pdf",
                            filename=out_name,
                            background=BackgroundTask(cleanup, inp, out))
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)

        raise safe_error(e, "edit-pdf")


# ─── 9. CMYK Color Converter ───────────────────────────────────
# Converts RGB PDF to CMYK color space for professional printing.
# Uses Ghostscript (already required for compress-pdf).

@app.post("/convert/cmyk")
@limiter.limit("10/minute")
async def convert_to_cmyk(request: Request, file: UploadFile = File(...)):
    """
    Convert an RGB PDF to CMYK color space using Ghostscript.
    Required for professional offset printing (ISO Coated v2 compatible).
    """
    inp = temp_path(".pdf")
    out = temp_path("_cmyk.pdf")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "convert-cmyk")

        gs_cmd = [
            "gs",
            "-dBATCH", "-dNOPAUSE", "-dQUIET",
            "-dSAFER",
            "-sDEVICE=pdfwrite",
            "-dCompatibilityLevel=1.4",
            "-sColorConversionStrategy=CMYK",
            "-dProcessColorModel=/DeviceCMYK",
            "-dOverrideICC=true",
            f"-sOutputFile={out}",
            str(inp)
        ]

        try:
            result = subprocess.run(
                gs_cmd, check=True,
                stdout=subprocess.DEVNULL,
                stderr=subprocess.PIPE,
                timeout=120
            )
        except subprocess.CalledProcessError as e:
            logger.error(f"[CMYK] Ghostscript failed: {e.stderr[:300] if e.stderr else 'no stderr'}")
            raise RuntimeError("Ghostscript CMYK conversion failed.")
        except FileNotFoundError:
            raise HTTPException(500, "Ghostscript is not installed on this server. Cannot perform CMYK conversion.")

        if not out.exists() or out.stat().st_size == 0:
            raise RuntimeError("CMYK output PDF not generated.")

        stem = safe_filename(file.filename, "document")
        out_name = Path(stem).stem + "_CMYK.pdf"
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "convert-cmyk")


# ─── 10. PDF Repair (Corrupt PDF Recovery) ─────────────────────
# Multi-strategy: Ghostscript → mutool → PyMuPDF open(repair=True)

@app.post("/repair-pdf")
@limiter.limit("10/minute")
async def repair_pdf(request: Request, file: UploadFile = File(...)):
    """
    Attempt to repair a corrupted or partially downloaded PDF.
    Tries 3 fallback strategies in order:
      1. Ghostscript (best for structural XREF/stream errors)
      2. mutool clean (MuPDF CLI — good for linearization errors)
      3. PyMuPDF open with garbage collection (last resort)
    """
    inp = temp_path(".pdf")
    out = temp_path("_repaired.pdf")
    try:
        raw_data = await safe_read_upload(file)
        inp.write_bytes(raw_data)

        repaired = False
        strategy_used = "unknown"

        # Strategy 1: Ghostscript rebuild
        try:
            gs_cmd = [
                "gs",
                "-dBATCH", "-dNOPAUSE", "-dQUIET",
                "-dSAFER",
                "-sDEVICE=pdfwrite",
                "-dCompatibilityLevel=1.4",
                f"-sOutputFile={out}",
                str(inp)
            ]
            r = subprocess.run(gs_cmd, capture_output=True, timeout=120)
            if r.returncode == 0 and out.exists() and out.stat().st_size > 0:
                repaired = True
                strategy_used = "Ghostscript (XREF rebuild)"
        except Exception as e:
            logger.warning(f"[REPAIR] Ghostscript failed: {type(e).__name__}")

        # Strategy 2: mutool clean
        if not repaired:
            try:
                mutool_path = shutil.which("mutool")
                if mutool_path:
                    r2 = subprocess.run(
                        [mutool_path, "clean", "-g", str(inp), str(out)],
                        capture_output=True, timeout=60
                    )
                    if r2.returncode == 0 and out.exists() and out.stat().st_size > 0:
                        repaired = True
                        strategy_used = "mutool (MuPDF clean)"
            except Exception as e:
                logger.warning(f"[REPAIR] mutool failed: {type(e).__name__}")

        # Strategy 3: PyMuPDF garbage collection
        if not repaired:
            try:
                import fitz
                doc = fitz.open(str(inp))
                doc.save(str(out), garbage=4, deflate=True, clean=True)
                doc.close()
                if out.exists() and out.stat().st_size > 0:
                    repaired = True
                    strategy_used = "PyMuPDF (garbage collection)"
            except Exception as e:
                logger.warning(f"[REPAIR] PyMuPDF failed: {type(e).__name__}")

        if not repaired:
            raise HTTPException(
                400,
                "This PDF is too severely damaged to recover. None of our repair strategies could reconstruct it."
            )

        stem = safe_filename(file.filename, "document")
        out_name = "repaired_" + stem
        response = FileResponse(
            str(out), media_type="application/pdf",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
        response.headers["X-Repair-Strategy"] = strategy_used
        response.headers["Access-Control-Expose-Headers"] = "X-Repair-Strategy"
        return response
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "repair-pdf")


# ─── 11. Webpage to PDF (Playwright) ───────────────────────────
# Converts any public URL to a high-quality paginated PDF.
# Requires: pip install playwright && playwright install chromium

WEBPAGE_TO_PDF_SEMAPHORE = asyncio.Semaphore(int(os.environ.get("MAX_CONCURRENT_PLAYWRIGHT", "2")))
ALLOWED_URL_SCHEMES = {"http", "https"}


def _validate_url(url: str) -> str:
    """Basic URL validation — only allow http/https, reject localhost/private IPs."""
    from urllib.parse import urlparse
    import ipaddress
    parsed = urlparse(url)
    if parsed.scheme not in ALLOWED_URL_SCHEMES:
        raise HTTPException(400, "Only http:// and https:// URLs are allowed.")
    hostname = parsed.hostname or ""
    # Block SSRF targets
    try:
        ip = ipaddress.ip_address(hostname)
        if ip.is_private or ip.is_loopback or ip.is_link_local:
            raise HTTPException(400, "Private/internal IP addresses are not allowed.")
    except ValueError:
        pass  # Not an IP — that's fine
    if hostname in ("localhost", "127.0.0.1", "::1"):
        raise HTTPException(400, "localhost is not allowed.")
    return url


@app.post("/convert/webpage-to-pdf")
@limiter.limit("5/minute")
async def webpage_to_pdf(
    request: Request,
    url: str = Form(...),
    paper_format: str = Form("A4"),
    landscape: bool = Form(False),
    margin: str = Form("1cm"),
    background: bool = Form(True),
):
    """
    Convert a public webpage URL to a paginated PDF using Playwright (Headless Chromium).
    """
    _validate_url(url)

    ALLOWED_FORMATS = {"A4", "A3", "Letter", "Legal", "Tabloid"}
    if paper_format not in ALLOWED_FORMATS:
        paper_format = "A4"

    out = temp_path(".pdf")
    try:
        try:
            from playwright.async_api import async_playwright
        except ImportError:
            raise HTTPException(
                500,
                "Playwright is not installed. Run: pip install playwright && playwright install chromium"
            )

        async with WEBPAGE_TO_PDF_SEMAPHORE:
            async with async_playwright() as p:
                browser = await p.chromium.launch(
                    headless=True,
                    args=["--no-sandbox", "--disable-dev-shm-usage"]
                )
                page = await browser.new_page()
                await page.goto(url, wait_until="networkidle", timeout=30000)
                await page.pdf(
                    path=str(out),
                    format=paper_format,
                    landscape=landscape,
                    print_background=background,
                    margin={
                        "top": margin, "right": margin,
                        "bottom": margin, "left": margin
                    }
                )
                await browser.close()

        if not out.exists() or out.stat().st_size == 0:
            raise RuntimeError("Playwright did not generate a PDF.")

        from urllib.parse import urlparse
        domain = urlparse(url).netloc.replace(".", "_")[:50]
        out_name = f"webpage_{domain}.pdf"
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=out_name,
            background=BackgroundTask(cleanup, out)
        )
    except HTTPException:
        cleanup(out)
        raise
    except Exception as e:
        cleanup(out)
        raise safe_error(e, "webpage-to-pdf")


# ─── 12. PDF Font Replacer ─────────────────────────────────────
# Extracts text from PDF pages and re-renders with chosen standard font.

ALLOWED_REPLACE_FONTS = {
    "helv": "Helvetica", "tiro": "Times-Roman",
    "cour": "Courier", "zadb": "ZapfDingbats"
}


@app.post("/convert/replace-font")
@limiter.limit("8/minute")
async def replace_font(
    request: Request,
    file: UploadFile = File(...),
    font_name: str = Form("helv"),
):
    """
    Replace all text in a PDF with a chosen standard font (Helvetica, Times, Courier).
    Text content and positions are preserved; only the font face changes.
    """
    if font_name not in ALLOWED_REPLACE_FONTS:
        font_name = "helv"

    import fitz
    import subprocess
    import shutil
    
    inp = temp_path(".pdf")
    out = temp_path("_font_replaced.pdf")
    no_text = temp_path("_notext.pdf")
    
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "replace-font")

        # 1. Use Ghostscript to strip all existing text (preserves vectors/images)
        gs_exe = "gs" if shutil.which("gs") else "gswin64c" if shutil.which("gswin64c") else "gswin32c"
        try:
            subprocess.run([
                gs_exe, "-q", "-dNOPAUSE", "-dBATCH", 
                "-sDEVICE=pdfwrite", "-dFILTERTEXT", 
                f"-sOutputFile={no_text}", str(inp)
            ], check=True, capture_output=True)
            has_gs = True
        except (subprocess.CalledProcessError, FileNotFoundError):
            has_gs = False

        src_doc = fitz.open(str(inp))
        if has_gs and no_text.exists():
            out_doc = fitz.open(str(no_text))
        else:
            out_doc = fitz.open()

        for page_num in range(len(src_doc)):
            src_page = src_doc[page_num]
            
            if has_gs and no_text.exists():
                new_page = out_doc[page_num]
            else:
                w, h = src_page.rect.width, src_page.rect.height
                new_page = out_doc.new_page(width=w, height=h)
                # Fallback: Redact old text (will leave white boxes, but better than double text)
                new_page.show_pdf_page(new_page.rect, src_doc, page_num)
                for block in src_page.get_text("dict")["blocks"]:
                    if block.get("type") == 0:
                        for line in block.get("lines", []):
                            for span in line.get("spans", []):
                                new_page.add_redact_annot(span["bbox"])
                new_page.apply_redactions()

            # 2. Re-insert text spans with new font
            for block in src_page.get_text("dict")["blocks"]:
                if block.get("type") != 0:
                    continue
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        txt = span.get("text", "").strip()
                        if not txt:
                            continue
                        x0, y0, _, y1 = span["bbox"]
                        size = max(4, span.get("size", 11))
                        
                        raw_color = span.get("color", 0)
                        if isinstance(raw_color, int):
                            r = ((raw_color >> 16) & 0xFF) / 255.0
                            g = ((raw_color >> 8) & 0xFF) / 255.0
                            b = (raw_color & 0xFF) / 255.0
                            color = (r, g, b)
                        else:
                            color = (0, 0, 0)
                            
                        new_page.insert_text(
                            fitz.Point(x0, y1 - 1),
                            txt,
                            fontname=font_name,
                            fontsize=size,
                            color=color,
                        )

        out_doc.save(str(out), garbage=3, deflate=True)
        src_doc.close()
        out_doc.close()

        font_label = ALLOWED_REPLACE_FONTS[font_name]
        stem = safe_filename(file.filename, "document")
        out_name = f"{Path(stem).stem}_{font_label}.pdf"
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "replace-font")


# ─── 13. Auto-Crop White Margins ───────────────────────────────
# Detects content bounding box per page and sets CropBox accordingly.

@app.post("/convert/auto-crop")
@limiter.limit("10/minute")
async def auto_crop_margins(
    request: Request,
    file: UploadFile = File(...),
    padding_pt: int = Form(10),
):
    """
    Automatically detect and remove white margins from every PDF page.
    Sets the CropBox to the content bounding box + optional padding.
    """
    import fitz
    padding_pt = max(0, min(padding_pt, 72))  # clamp 0–72pt

    inp = temp_path(".pdf")
    out = temp_path("_autocropped.pdf")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "auto-crop")

        doc = fitz.open(str(inp))
        for page in doc:
            # get_text("blocks") returns (x0, y0, x1, y1, text, ...) per block
            blocks = page.get_text("blocks")
            # Also include drawings/images
            paths = page.get_drawings()
            images = page.get_image_rects(full=True)  # Returns list of (Rect, transform, ...)

            all_rects = []
            for b in blocks:
                if b[4].strip():  # has text
                    all_rects.append(fitz.Rect(b[:4]))
            for path in paths:
                if path.get("rect"):
                    all_rects.append(path["rect"])
            for img_info in images:
                if isinstance(img_info, (list, tuple)) and len(img_info) >= 1:
                    rect = img_info[0] if isinstance(img_info[0], fitz.Rect) else None
                    if rect:
                        all_rects.append(rect)

            if all_rects:
                content_bbox = all_rects[0]
                for r in all_rects[1:]:
                    content_bbox = content_bbox | r  # union
                # Add padding
                padded = fitz.Rect(
                    max(0, content_bbox.x0 - padding_pt),
                    max(0, content_bbox.y0 - padding_pt),
                    min(page.rect.width, content_bbox.x1 + padding_pt),
                    min(page.rect.height, content_bbox.y1 + padding_pt),
                )
                page.set_cropbox(padded)

        doc.save(str(out), garbage=3, deflate=True)
        doc.close()

        stem = safe_filename(file.filename, "document")
        out_name = "autocropped_" + stem
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "auto-crop")


# ─── 14. PDF Link Extractor ─────────────────────────────────────

@app.post("/analyze/links")
@limiter.limit("15/minute")
async def extract_links(request: Request, file: UploadFile = File(...)):
    """
    Extract all hyperlinks from a PDF and return them as JSON.
    """
    import fitz
    inp = temp_path(".pdf")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "extract-links")

        doc = fitz.open(str(inp))
        links_found = []
        for i, page in enumerate(doc):
            for link in page.get_links():
                uri = link.get("uri", "")
                if uri:
                    links_found.append({
                        "page": i + 1,
                        "url": uri,
                        "rect": [round(v, 1) for v in link.get("from", fitz.Rect())]
                    })
        doc.close()

        return JSONResponse({
            "success": True,
            "total_links": len(links_found),
            "links": links_found
        })
    except HTTPException:
        raise
    except Exception as e:
        raise safe_error(e, "extract-links")
    finally:
        cleanup(inp)


# ─── 15. Compress PDF ───────────────────────────────────────────
# quality: screen (max), ebook (good), printer (high quality), prepress (best)

ALLOWED_QUALITY_VALUES = {"screen", "ebook", "printer", "prepress"}

@app.post("/compress-pdf")
@limiter.limit("15/minute")
async def compress_pdf(
    request: Request,
    file: UploadFile = File(...),
    quality: str = Form("ebook")   # screen | ebook | printer | prepress
):
    if quality not in ALLOWED_QUALITY_VALUES:
        quality = "ebook"

    inp = temp_path(".pdf")
    out = temp_path("_compressed.pdf")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "compress-pdf")

        gs_cmd = [
            "gs", "-sDEVICE=pdfwrite", "-dCompatibilityLevel=1.4",
            f"-dPDFSETTINGS=/{quality}", "-dNOPAUSE", "-dQUIET",
            "-dBATCH", f"-sOutputFile={out}", str(inp)
        ]
        try:
            subprocess.run(gs_cmd, check=True, stdout=subprocess.DEVNULL,
                           stderr=subprocess.DEVNULL, timeout=120)
        except Exception:
            import fitz
            doc = fitz.open(inp)
            doc.save(str(out), garbage=4, deflate=True, clean=True)
            doc.close()

        stem = safe_filename(file.filename, "document")
        out_name = "compressed_" + stem
        return FileResponse(str(out), media_type="application/pdf",
                            filename=out_name,
                            background=BackgroundTask(cleanup, inp, out))
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "compress-pdf")


# ─── 10. Extract Tables ────────────────────────────────────────

@app.post("/extract-tables")
@limiter.limit("10/minute")
async def extract_tables(request: Request, file: UploadFile = File(...)):
    try:
        import camelot, pandas as pd
    except ImportError:
        raise HTTPException(500, "Run: pip install camelot-py[cv] pandas openpyxl")

    inp = temp_path(".pdf")
    out = temp_path(".xlsx")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "extract-tables")

        tables = None
        for flavor in ["lattice", "stream"]:
            try:
                t = camelot.read_pdf(str(inp), pages="all", flavor=flavor)
                if t and len(t) > 0:
                    tables = t
                    break
            except Exception:
                continue

        if not tables:
            raise HTTPException(400, "No tables found in this PDF.")

        with pd.ExcelWriter(str(out), engine="openpyxl") as writer:
            for i, table in enumerate(tables):
                sheet_name = f"Table_{i+1}"[:31]
                df = table.df
                # Use first row as header if it looks like one
                df.to_excel(writer, sheet_name=sheet_name, index=False, header=False)

        stem = safe_filename(file.filename, "document")
        out_name = "tables_" + Path(stem).stem + ".xlsx"
        return FileResponse(
            str(out),
            media_type="application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
            filename=out_name,
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "extract-tables")


# ─── 11. Extract Images ────────────────────────────────────────

@app.post("/extract-images")
@limiter.limit("10/minute")
async def extract_images(request: Request, file: UploadFile = File(...)):
    import zipfile
    inp = temp_path(".pdf")
    out_zip = temp_path(".zip")
    img_dir = temp_path("_imgdir")
    img_dir.mkdir(exist_ok=True)

    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "extract-images")

        import fitz
        doc = fitz.open(inp)
        count = 0
        seen_xrefs = set()  # deduplicate images appearing on multiple pages

        for page_num in range(len(doc)):
            page = doc[page_num]
            for img in page.get_images(full=True):
                xref = img[0]
                if xref in seen_xrefs:
                    continue
                seen_xrefs.add(xref)
                try:
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    if len(image_bytes) < 500:  # skip tiny/empty images
                        continue
                    ext = base_image.get("ext", "png")
                    count += 1
                    img_filepath = img_dir / f"page_{page_num+1}_img_{count}.{ext}"
                    img_filepath.write_bytes(image_bytes)
                except Exception:
                    continue

        doc.close()

        if count == 0:
            cleanup(inp, out_zip, img_dir)
            raise HTTPException(400, "No images found in this PDF.")

        with zipfile.ZipFile(str(out_zip), "w", zipfile.ZIP_DEFLATED) as zipf:
            for f in img_dir.iterdir():
                zipf.write(f, f.name)

        stem = safe_filename(file.filename, "document")
        out_name = "images_" + Path(stem).stem + ".zip"
        return FileResponse(str(out_zip), media_type="application/zip",
                            filename=out_name,
                            background=BackgroundTask(cleanup, inp, out_zip, img_dir))
    except HTTPException:
        cleanup(inp, out_zip, img_dir)
        raise
    except Exception as e:
        cleanup(inp, out_zip, img_dir)
        raise safe_error(e, "extract-images")


# ─── 12. Search & Replace ──────────────────────────────────────
# Matches the original font size so replaced text looks correct.

@app.post("/search-replace")
@limiter.limit("15/minute")
async def search_replace(
    request: Request,
    file: UploadFile = File(...),
    search_term: str = Form(...),
    replace_term: str = Form(...)
):
    inp = temp_path(".pdf")
    out = temp_path("_replaced.pdf")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "search-replace")

        import fitz
        doc = fitz.open(inp)
        match_count = 0

        for page in doc:
            instances = page.search_for(search_term)
            if not instances:
                continue
            match_count += len(instances)

            # Collect original font info before redacting
            font_info = {}
            blocks = page.get_text("dict").get("blocks", [])
            for block in blocks:
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        if search_term.lower() in span["text"].lower():
                            font_info = {
                                "size": span.get("size", 11),
                                "color": span.get("color", 0)
                            }

            # Redact (erase old text with white fill)
            for inst in instances:
                page.add_redact_annot(inst, fill=(1, 1, 1))
            page.apply_redactions()

            # Insert replacement with matched font size
            font_size = font_info.get("size", 11)
            # Convert packed int color to tuple
            raw_color = font_info.get("color", 0)
            if isinstance(raw_color, int):
                r = ((raw_color >> 16) & 0xFF) / 255.0
                g = ((raw_color >> 8) & 0xFF) / 255.0
                b = (raw_color & 0xFF) / 255.0
                text_color = (r, g, b)
            else:
                text_color = (0, 0, 0)

            # Determine font for replacement text
            _sr_fonts_dir = Path(__file__).parent.parent / "public" / "fonts"
            _sr_fontname = "helv"
            if has_non_latin(replace_term):
                _sr_font_file = _sr_fonts_dir / "NotoSansDevanagari-Regular.ttf"
                if not _sr_font_file.exists():
                    _sr_font_file = Path(__file__).parent / "fonts" / "NotoSansDevanagari-Regular.ttf"
                if _sr_font_file.exists():
                    _sr_fontname = "NotoDevaRep"
                    try:
                        page.insert_font(fontname=_sr_fontname, fontbuffer=_sr_font_file.read_bytes())
                    except Exception:
                        pass  # Already registered on this page

            for inst in instances:
                page.insert_text(
                    fitz.Point(inst.x0, inst.y1 - 1),
                    replace_term,
                    fontname=_sr_fontname,
                    fontsize=font_size,
                    color=text_color
                )

        doc.save(str(out))
        doc.close()

        stem = safe_filename(file.filename, "document")
        out_name = "replaced_" + stem
        return FileResponse(
            str(out), media_type="application/pdf",
            filename=out_name,
            headers={
                "X-Match-Count": str(match_count),
                "Access-Control-Expose-Headers": "X-Match-Count"
            },
            background=BackgroundTask(cleanup, inp, out)
        )
    except HTTPException:
        cleanup(inp, out)
        raise
    except Exception as e:
        cleanup(inp, out)
        raise safe_error(e, "search-replace")


# ─── 13. OCR Analyze (Detect Scan + Return Text) ──────────────

@app.post("/ocr/analyze")
@limiter.limit("10/minute")
async def ocr_analyze(request: Request, file: UploadFile = File(...)):
    """
    Analyze a PDF: detect if it's scanned, extract text per page.
    Returns JSON with is_scanned flag and text for each page.
    """
    inp = temp_path(".pdf")
    try:
        inp.write_bytes(await safe_read_upload(file))
        check_pdf_page_count(inp, "ocr-analyze")

        import fitz
        doc = fitz.open(inp)
        pages_data = []
        total_chars = 0
        for i, page in enumerate(doc):
            text = page.get_text("text")
            total_chars += len(text.strip())
            pages_data.append({"page": i + 1, "char_count": len(text.strip()), "text_preview": text[:300]})
        doc.close()

        scanned = is_pdf_scanned(inp)
        return JSONResponse({
            "success": True,
            "is_scanned": scanned,
            "total_chars": total_chars,
            "pages": pages_data
        })
    except HTTPException:
        raise
    except Exception as e:
        raise safe_error(e, "ocr-analyze")
    finally:
        cleanup(inp)
